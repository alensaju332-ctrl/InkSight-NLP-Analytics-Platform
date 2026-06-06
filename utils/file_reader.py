import os
import logging

logger = logging.getLogger(__name__)

# Optional dependencies: try importing them and provide safe fallbacks if unavailable
try:
    import pdfplumber
except Exception:
    pdfplumber = None

try:
    import PyPDF2
except Exception:
    PyPDF2 = None

try:
    from docx import Document
except Exception:
    Document = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

import csv

try:
    import markdown
except Exception:
    markdown = None

try:
    from odf import text, teletype
    from odf.opendocument import load
except Exception:
    text = None
    teletype = None
    load = None

try:
    import openpyxl
except Exception:
    openpyxl = None

try:
    import chardet
except Exception:
    chardet = None

class FileReader:
    """Universal file reader for multiple formats"""
    
    def __init__(self):
        self.supported_formats = {
            '.txt': self.read_txt,
            '.docx': self.read_docx,
            '.doc': self.read_doc,
            '.pdf': self.read_pdf,
            '.ppt': self.read_ppt,
            '.pptx': self.read_pptx,
            '.odt': self.read_odt,
            '.rtf': self.read_rtf,
            '.md': self.read_markdown,
            '.csv': self.read_csv,
            '.xlsx': self.read_xlsx,
            '.xls': self.read_xlsx
        }
    
    def read_file(self, filepath):
        """Read file based on extension"""
        ext = os.path.splitext(filepath)[1].lower()
        
        if ext in self.supported_formats:
            try:
                return self.supported_formats[ext](filepath)
            except Exception as e:
                logger.error(f"Error reading {ext} file: {str(e)}")
                # Try with basic text reading as fallback
                return self.read_txt_with_encoding(filepath)
        else:
            raise ValueError(f"Unsupported file format: {ext}")
    
    def read_txt(self, filepath):
        """Read plain text file with encoding detection"""
        return self.read_txt_with_encoding(filepath)
    
    def read_txt_with_encoding(self, filepath):
        """Read text file with automatic encoding detection"""
        # If chardet is available, use it to detect encoding; otherwise try common encodings
        if chardet is not None:
            with open(filepath, 'rb') as f:
                raw_data = f.read()
                try:
                    result = chardet.detect(raw_data)
                    encoding = result.get('encoding') or 'utf-8'
                except Exception:
                    encoding = 'utf-8'
            # Read with detected encoding
            try:
                with open(filepath, 'r', encoding=encoding) as f:
                    return f.read()
            except Exception:
                # Fallback to utf-8 with ignore
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
        else:
            # chardet not available — try common encodings
            encodings_to_try = ['utf-8', 'utf-8-sig', 'cp1252', 'latin-1']
            for enc in encodings_to_try:
                try:
                    with open(filepath, 'r', encoding=enc) as f:
                        return f.read()
                except Exception:
                    continue
            # Final fallback: read as utf-8 ignoring errors
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
    
    def read_docx(self, filepath):
        """Read DOCX file"""
        if Document is None:
            raise ImportError("python-docx is not available. Install it to read DOCX files.")
        
        doc = Document(filepath)
        text = []
        
        # Extract text from paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text.append(paragraph.text)
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text.append(' | '.join(row_text))
        
        return '\n'.join(text)
    
    def read_doc(self, filepath):
        """Read DOC file (legacy format)"""
        # For .doc files, we'll use python-docx2txt or convert to docx
        # As a fallback, try reading as docx
        try:
            return self.read_docx(filepath)
        except:
            # If pypandoc is available, use it
            try:
                import pypandoc
                return pypandoc.convert_file(filepath, 'plain')
            except:
                raise ValueError("Cannot read .doc file. Please convert to .docx format.")
    
    def read_pdf(self, filepath):
        """Read PDF file using multiple methods for better extraction"""
        import re
        text = []

        # Method 1: Try pdfplumber first (better for complex PDFs)
        if pdfplumber is not None:
            try:
                with pdfplumber.open(filepath) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text.append(page_text)

                if text:
                    extracted_text = '\n'.join(text)
                    # Clean PDF artifacts
                    return self._clean_pdf_text(extracted_text)
            except Exception as e:
                logger.warning(f"pdfplumber failed: {str(e)}")

        # Method 2: Fallback to PyPDF2
        if PyPDF2 is not None:
            try:
                with open(filepath, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    for page_num in range(len(pdf_reader.pages)):
                        page = pdf_reader.pages[page_num]
                        page_text = page.extract_text()
                        if page_text.strip():
                            text.append(page_text)

                extracted_text = '\n'.join(text)
                # Clean PDF artifacts
                return self._clean_pdf_text(extracted_text)
            except Exception as e:
                logger.error(f"PyPDF2 failed: {str(e)}")
                raise ValueError("Unable to extract text from PDF")

        raise ImportError("No suitable PDF library found. Install pdfplumber or PyPDF2.")

    def _clean_pdf_text(self, text):
        """Clean PDF extraction artifacts"""
        import re

        # Fix 1: Remove line-break hyphens (e.g., "deci-\nsionmaking" → "decisionmaking")
        # This fixes the word count issue where PDFs show 101 words instead of 100
        text = re.sub(r'-\s*\n\s*', '', text)

        # Fix 2: Replace multiple spaces with single space
        text = re.sub(r' {2,}', ' ', text)

        # Fix 3: Remove form feed characters (page breaks)
        text = text.replace('\x0c', '\n')

        # Fix 4: Normalize line endings
        text = re.sub(r'\r\n', '\n', text)

        return text.strip()
    
    def read_pptx(self, filepath):
        """Read PowerPoint PPTX file"""
        if Presentation is None:
            raise ImportError("python-pptx is not available. Install it to read PPTX files.")
        
        prs = Presentation(filepath)
        text = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"Slide {slide_num}:"]
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
                
                # Extract text from tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            slide_text.append(' | '.join(row_text))
            
            if len(slide_text) > 1:  # If there's more than just the slide number
                text.append('\n'.join(slide_text))
        
        return '\n\n'.join(text)
    
    def read_ppt(self, filepath):
        """Read PPT file (legacy format)"""
        # Try to read as pptx first
        try:
            return self.read_pptx(filepath)
        except:
            # If pypandoc is available, use it
            try:
                import pypandoc
                return pypandoc.convert_file(filepath, 'plain')
            except:
                raise ValueError("Cannot read .ppt file. Please convert to .pptx format.")
    
    def read_odt(self, filepath):
        """Read OpenDocument Text file"""
        if text is None or load is None:
            raise ImportError("python-odf is not available. Install it to read ODT files.")
        
        doc = load(filepath)
        text_content = []
        
        for element in doc.getElementsByType(text.P):
            paragraph_text = teletype.extractText(element)
            if paragraph_text.strip():
                text_content.append(paragraph_text)
        
        return '\n'.join(text_content)
    
    def read_rtf(self, filepath):
        """Read RTF file"""
        try:
            import pypandoc
            return pypandoc.convert_file(filepath, 'plain')
        except:
            # Basic RTF stripping (very simple, may not work for complex RTF)
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
                # Remove RTF commands (very basic)
                import re
                text = re.sub(r'\\[a-z]+\d*\s?', '', content)
                text = re.sub(r'[{}]', '', text)
                return text.strip()
    
    def read_markdown(self, filepath):
        """Read Markdown file and convert to plain text"""
        if markdown is None:
            raise ImportError("markdown is not available. Install it to read Markdown files.")
        
        with open(filepath, 'r', encoding='utf-8') as f:
            md_content = f.read()
        
        # Convert markdown to HTML then strip tags
        html = markdown.markdown(md_content)
        # Remove HTML tags
        import re
        text = re.sub(r'<[^>]+>', '', html)
        return text
    
    def read_csv(self, filepath):
        """Read CSV file"""
        text = []
        
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            csv_reader = csv.reader(f)
            for row in csv_reader:
                if row:  # Skip empty rows
                    text.append(' '.join(str(cell) for cell in row if cell))
        
        return '\n'.join(text)
    
    def read_xlsx(self, filepath):
        """Read Excel file"""
        if openpyxl is None:
            raise ImportError("openpyxl is not available. Install it to read XLSX files.")
        
        workbook = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
        text = []
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_text = [f"Sheet: {sheet_name}"]
            
            for row in sheet.iter_rows(values_only=True):
                row_text = []
                for cell in row:
                    if cell is not None and str(cell).strip():
                        row_text.append(str(cell))
                if row_text:
                    sheet_text.append(' | '.join(row_text))
            
            if len(sheet_text) > 1:  # If there's more than just the sheet name
                text.append('\n'.join(sheet_text))
        
        workbook.close()
        return '\n\n'.join(text)